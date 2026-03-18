#!/usr/bin/env python3
"""
Build the GSoC proposal presentation in Metaflow's visual style.

Style reference: Valay Dave's PyData NYC 2023 presentation
- Background: warm beige (#F0ECE3)
- Text: dark charcoal (#2D2D2D)
- Accents: Metaflow purple (#7B68AE), Metaflow blue (#1A56DB)
- Font: clean sans-serif (Calibri as fallback)
- Lots of whitespace, bold titles, minimal bullets
- Memes and humor where appropriate
- Code in dark blocks
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (Metaflow palette) ──
BG = RGBColor(0xF0, 0xEC, 0xE3)       # warm beige
TEXT = RGBColor(0x2D, 0x2D, 0x2D)      # dark charcoal
TEXT_LIGHT = RGBColor(0x6B, 0x6B, 0x6B) # grey for subtitles
PURPLE = RGBColor(0x7B, 0x68, 0xAE)    # Metaflow purple
BLUE = RGBColor(0x1A, 0x56, 0xDB)      # Metaflow blue
RED = RGBColor(0xE8, 0x57, 0x57)       # red for "bad" numbers
GREEN = RGBColor(0x4A, 0xA8, 0x6B)     # green for "good" numbers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x2D, 0x2D, 0x2D)   # dark code background
LIGHT_PURPLE = RGBColor(0xE8, 0xE0, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18,
             bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=16,
                  color=TEXT, line_spacing=1.5, bold_first=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(size * (line_spacing - 1))
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_code_block(slide, left, top, width, height, code, size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        p.font.name = "Courier New"
    return shape

def add_pill(slide, left, top, width, text, color=PURPLE, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, left, top, width):
    shape = slide.shapes.add_connector(1, Inches(left), Inches(top),
                                        Inches(left + width), Inches(top))
    shape.line.color.rgb = TEXT_LIGHT
    shape.line.width = Pt(1)
    return shape

# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_text(slide, 0.8, 1.0, 8, 2.5,
         "Agent-Friendly\nMetaflow Client",
         size=48, bold=True)
add_text(slide, 0.8, 3.5, 6, 0.5,
         "Analyzing and Addressing Client API Inefficiencies",
         size=20, color=TEXT_LIGHT)
add_text(slide, 0.8, 4.5, 4, 0.4, "landigf", size=18, bold=True)
add_text(slide, 0.8, 5.0, 4, 0.4,
         "MSc Computer Science, ETH Zurich", size=14, color=TEXT_LIGHT)
add_text(slide, 0.8, 5.4, 4, 0.4,
         "GSoC 2026 Proposal | Mentor: Valay Dave", size=14, color=TEXT_LIGHT)
add_text(slide, 10.5, 6.8, 2.5, 0.4, "Outerbounds",
         size=14, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: Outline
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.8, 4, 0.8, "Outline", size=36, bold=True)

sections = [
    ("01", "The Problem: Why agents struggle with Metaflow today"),
    ("02", "The Discovery: Infrastructure that already exists"),
    ("03", "The Benchmark: Three paths, same question, real numbers"),
    ("04", "The Solution: What we're building"),
    ("05", "GSoC Timeline: 350 hours, four phases"),
]
for i, (num, title) in enumerate(sections):
    y = 2.0 + i * 1.0
    add_text(slide, 5.5, y - 0.15, 0.6, 0.4, num,
             size=13, color=PURPLE, bold=True)
    add_line(slide, 6.1, y, 6.5)
    add_text(slide, 6.1, y + 0.05, 6.5, 0.5, title, size=18)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: Section divider - The Problem
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1.5, 2.8, 10, 2,
         "What happens when an agent\nasks a simple question?",
         size=40, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: The question
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.6, 8, 0.6, "The agent asks:", size=24, bold=True)
add_text(slide, 0.8, 1.3, 10, 0.8,
         '"Which tasks failed in my latest ForeachFlow run?"',
         size=28, bold=True, color=PURPLE)

add_text(slide, 0.8, 2.5, 5, 0.5, "What the Client API does:", size=18, bold=True)

add_code_block(slide, 0.8, 3.1, 5.5, 3.5, """for step in run:
    for task in step:
        if not task.successful:
            # Inside task.successful:
            # 1. HTTP GET /artifacts/_success
            # 2. Read boolean from S3/local
            # 3. Unpickle Python object
            # 4. Return True/False
            #
            # PER TASK. 50 tasks = 50x this.""", size=13)

add_text(slide, 7.5, 2.5, 5, 0.5, "The cost:", size=18, bold=True)

# Big numbers
add_text(slide, 7.5, 3.2, 2.5, 0.7, "56", size=60, bold=True, color=RED)
add_text(slide, 9.5, 3.5, 3, 0.5, "HTTP calls", size=22, color=TEXT_LIGHT)

add_text(slide, 7.5, 4.3, 2.5, 0.7, "3.8s", size=60, bold=True, color=RED)
add_text(slide, 9.8, 4.6, 3, 0.5, "wall-clock time", size=22, color=TEXT_LIGHT)

add_text(slide, 7.5, 5.5, 2.5, 0.7, "50", size=60, bold=True, color=RED)
add_text(slide, 9.5, 5.8, 3, 0.5, "datastore reads\nto unpickle 50 booleans",
         size=18, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: Architecture diagram
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 10, 0.6,
         "Two services, same database, different capabilities",
         size=28, bold=True)

# PostgreSQL box
add_rect(slide, 3.5, 1.5, 6, 0.8, PURPLE)
add_text(slide, 3.5, 1.55, 6, 0.7, "PostgreSQL",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Metadata Service
add_rect(slide, 1.5, 3.0, 4, 2.5, RGBColor(0xE8, 0xE8, 0xE8))
add_text(slide, 1.7, 3.1, 3.6, 0.4, "Metadata Service (port 8080)",
         size=16, bold=True)
add_multiline(slide, 1.7, 3.6, 3.6, 2.0, [
    "Every deployment has this",
    "",
    "No pagination",
    "No filtering",
    "No status field",
    "Returns ALL data",
], size=13, color=TEXT_LIGHT)

# UI Backend
add_rect(slide, 7.5, 3.0, 4, 2.5, RGBColor(0xE8, 0xE8, 0xE8))
add_text(slide, 7.7, 3.1, 3.6, 0.4, "UI Backend (port 8083)",
         size=16, bold=True)
add_multiline(slide, 7.7, 3.6, 3.6, 2.0, [
    "Optional, for Web UI only",
    "",
    "?_limit=10",
    "?status=failed",
    "?_order=-ts_epoch",
    "Status via SQL JOIN",
], size=13, color=GREEN)

# Bottom labels
add_text(slide, 1.5, 5.8, 4, 0.8,
         "Python Client API\n(what agents use today)",
         size=15, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 7.5, 5.8, 4, 0.8,
         "Metaflow Web UI\n(what humans see)",
         size=15, bold=True, align=PP_ALIGN.CENTER)

# Arrow annotation
add_text(slide, 3.0, 6.7, 7, 0.4,
         "The Client API only talks to the Metadata Service. It has no idea the UI Backend exists.",
         size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: Section divider - Discovery
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1.5, 2.8, 10, 2,
         "Three discoveries that\nchange everything",
         size=40, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: Discovery 1 - filter_tasks_by_metadata
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_pill(slide, 0.8, 0.6, 2.5, "Discovery 1")
add_text(slide, 0.8, 1.2, 11, 0.8,
         "filter_tasks_by_metadata already finds failures",
         size=30, bold=True)

add_text(slide, 0.8, 2.2, 11, 0.6,
         "Since metadata service v2.5.0, this endpoint exists but nobody uses it for failure detection:",
         size=16, color=TEXT_LIGHT)

add_code_block(slide, 0.8, 3.0, 11, 2.0, """# ONE HTTP call. Returns all failed task IDs.
ServiceMetadataProvider.filter_tasks_by_metadata(
    "ForeachFlow", "8", "process", "attempt_ok", "False"
)
# --> ["ForeachFlow/8/process/126"]

# Instead of iterating 50 tasks and fetching _success artifact for each one.""", size=15)

add_text(slide, 0.8, 5.3, 11, 1.5,
         "When a task finishes, the Metaflow runtime writes attempt_ok = \"True\" or \"False\" "
         "to the metadata_v3 table. This is the same field the UI Backend reads via SQL. "
         "The metadata service has this data. We just need to use it.",
         size=16, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: Discovery 2 - DB already has LIMIT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_pill(slide, 0.8, 0.6, 2.5, "Discovery 2")
add_text(slide, 0.8, 1.2, 11, 0.8,
         "The DB layer already supports LIMIT and ORDER BY",
         size=30, bold=True)

add_text(slide, 0.8, 2.2, 11, 0.6,
         "postgres_async_db.py has find_records() with limit, offset, and order parameters. "
         "The HTTP endpoints just don't expose them.",
         size=16, color=TEXT_LIGHT)

add_code_block(slide, 0.8, 3.0, 5.2, 2.5, """# What the metadata service does today:
SELECT * FROM runs_v3
WHERE flow_id = 'MyFlow'
-- returns ALL runs, unbounded

# What it COULD do (~15 lines of change):
SELECT * FROM runs_v3
WHERE flow_id = 'MyFlow'
ORDER BY ts_epoch DESC
LIMIT 10""", size=14)

add_text(slide, 6.5, 3.0, 5.5, 2.5,
         "The infrastructure is there.\n\n"
         "find_records() already accepts\n"
         "limit, offset, and order.\n\n"
         "~15 lines per endpoint to expose it.\n\n"
         "Benefits everyone, not just agents.",
         size=17)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: Discovery 3 - Status is simple
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_pill(slide, 0.8, 0.6, 2.5, "Discovery 3")
add_text(slide, 0.8, 1.2, 11, 0.8,
         "Task status computation is simpler than it looks",
         size=30, bold=True)

add_text(slide, 0.8, 2.2, 5.5, 0.5,
         "UI Backend: 300 lines, complex JOINs", size=16, color=RED, bold=True)
add_code_block(slide, 0.8, 2.8, 5.2, 2.8, """-- UI Backend task status (simplified)
CASE
  WHEN attempt_ok IS TRUE  -> 'completed'
  WHEN attempt_ok IS FALSE -> 'failed'
  WHEN heartbeat stale     -> 'failed'
  WHEN old run no hb       -> 'failed'
  WHEN no attempt          -> 'pending'
  ELSE                     -> 'running'
END
-- 100+ lines of lateral JOINs""", size=13)

add_text(slide, 6.8, 2.2, 5.5, 0.5,
         "What we need: 3 lines", size=16, color=GREEN, bold=True)
add_code_block(slide, 6.8, 2.8, 5.2, 2.0, """-- Simplified status (covers 90% of cases)
CASE
  WHEN attempt_ok = 'True'  -> 'completed'
  WHEN attempt_ok = 'False' -> 'failed'
  ELSE                      -> 'running'
END""", size=14)

add_text(slide, 6.8, 5.2, 5.5, 1.5,
         "The core signal is attempt_ok in the metadata table. "
         "The other 97 lines handle edge cases (stale heartbeats, old runs). "
         "For agent queries, the simplified version is enough.",
         size=15, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: Section divider - Benchmark
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1.5, 2.8, 10, 2,
         "The Benchmark\nThree paths, one question, real numbers",
         size=40, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: Three-path results
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         '"Which tasks failed?" — 50-task foreach, 1 failure',
         size=26, bold=True)

add_text(slide, 0.8, 1.3, 11, 0.4,
         "All three paths return the same answer. Same database. Same failed task.",
         size=16, color=TEXT_LIGHT)

# Path A - red bar
bar_top = 2.2
add_rect(slide, 0.8, bar_top, 10.0, 1.2, RGBColor(0xFD, 0xE8, 0xE8))
add_text(slide, 1.0, bar_top + 0.1, 3, 0.4,
         "A: Naive Client API", size=16, bold=True, color=RED)
add_text(slide, 1.0, bar_top + 0.55, 4, 0.4,
         "Iterate tasks, fetch _success artifact per task", size=13, color=TEXT_LIGHT)
add_text(slide, 7.5, bar_top + 0.15, 3, 0.9,
         "56 calls\n3,800ms", size=24, bold=True, color=RED, align=PP_ALIGN.RIGHT)

# Path C - green bar
bar_top = 3.7
add_rect(slide, 0.8, bar_top, 10.0, 1.2, RGBColor(0xE8, 0xF5, 0xE8))
add_text(slide, 1.0, bar_top + 0.1, 4, 0.4,
         "C: Smart Metadata", size=16, bold=True, color=GREEN)
add_text(slide, 1.0, bar_top + 0.55, 5, 0.4,
         "filter_tasks_by_metadata — no UI Backend needed", size=13, color=TEXT_LIGHT)
add_text(slide, 7.5, bar_top + 0.15, 3, 0.9,
         "4 calls\n349ms", size=24, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# Path B - blue bar
bar_top = 5.2
add_rect(slide, 0.8, bar_top, 10.0, 1.2, RGBColor(0xE8, 0xEC, 0xF5))
add_text(slide, 1.0, bar_top + 0.1, 4, 0.4,
         "B: UI Backend", size=16, bold=True, color=BLUE)
add_text(slide, 1.0, bar_top + 0.55, 5, 0.4,
         "Requires extra service (port 8083)", size=13, color=TEXT_LIGHT)
add_text(slide, 7.5, bar_top + 0.15, 3, 0.9,
         "2 calls\n35ms", size=24, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

# Takeaway
add_text(slide, 0.8, 6.6, 11, 0.5,
         "Smart Metadata: 8.6x faster than today, works everywhere, no extra infrastructure.",
         size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: DAG visualization placeholder
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "Metaflow benchmarks itself — visible in the UI",
         size=26, bold=True)
add_text(slide, 0.8, 1.2, 11, 0.5,
         "BenchmarkThreePaths flow: three branches, same question, join to compare",
         size=16, color=TEXT_LIGHT)

# DAG diagram
add_rect(slide, 5.5, 2.2, 2.0, 0.6, PURPLE)
add_text(slide, 5.5, 2.22, 2.0, 0.55, "start",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 1.5, 3.5, 3.0, 0.8, RED)
add_text(slide, 1.5, 3.5, 3.0, 0.4, "path_a_naive",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.9, 3.0, 0.4, "56 calls, 3.8s",
         size=12, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

add_rect(slide, 5.0, 3.5, 3.0, 0.8, GREEN)
add_text(slide, 5.0, 3.5, 3.0, 0.4, "path_c_smart_meta",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 5.0, 3.9, 3.0, 0.4, "4 calls, 349ms",
         size=12, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

add_rect(slide, 8.5, 3.5, 3.0, 0.8, BLUE)
add_text(slide, 8.5, 3.5, 3.0, 0.4, "path_b_ui_backend",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 8.5, 3.9, 3.0, 0.4, "2 calls, 35ms",
         size=12, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

add_rect(slide, 5.5, 5.0, 2.0, 0.6, PURPLE)
add_text(slide, 5.5, 5.02, 2.0, 0.55, "compare",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 5.5, 6.0, 2.0, 0.6, PURPLE)
add_text(slide, 5.5, 6.02, 2.0, 0.55, "end",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, 0.8, 6.8, 11, 0.4,
         "[Replace this slide with actual Metaflow UI screenshot of BenchmarkThreePaths/10]",
         size=14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: Full audit results
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "All 6 GSoC use cases — benchmarked",
         size=26, bold=True)

# Table header
y = 1.5
add_rect(slide, 0.8, y, 11.5, 0.5, PURPLE)
cols = [("UC", 0.5), ("Use Case", 4.5), ("Client API", 2.0),
        ("Smart Path", 2.0), ("Speedup", 1.5)]
x = 0.9
for label, w in cols:
    add_text(slide, x, y + 0.05, w, 0.4, label,
             size=13, bold=True, color=WHITE)
    x += w + 0.2

# Table rows
data = [
    ("1", "List runs by status", "4 calls, 134ms", "1 call, 20ms", "6.7x"),
    ("2", "Filter runs by time range", "2 calls, 18ms", "1 call, 12ms", "1.5x"),
    ("3", "Find failed tasks + errors", "57 calls, 2006ms", "5 calls, 105ms", "19.1x"),
    ("4", "Artifact metadata (no data)", "5 calls, 39ms", "4 calls, 133ms", "0.3x"),
    ("5", "Bounded log output", "6 calls, 93ms", "4 calls, 231ms", "0.4x"),
    ("6", "Cross-run artifact search", "8 calls, 93ms", "5 calls, 51ms", "1.8x"),
]
for i, row in enumerate(data):
    y = 2.2 + i * 0.7
    bg_color = RGBColor(0xF8, 0xF6, 0xF0) if i % 2 == 0 else BG
    add_rect(slide, 0.8, y, 11.5, 0.6, bg_color)
    x = 0.9
    for j, (_, w) in enumerate(cols):
        color = TEXT
        if j == 4:
            val = float(row[j].replace("x", ""))
            color = GREEN if val >= 1.0 else RED
        add_text(slide, x, y + 0.1, w, 0.4, row[j], size=13, color=color)
        x += w + 0.2

# Summary
add_text(slide, 0.8, 6.5, 11.5, 0.7,
         "UC3 (find failures) is the dominant bottleneck: 57 of 82 total calls (70%). "
         "UC4 and UC5 need different fixes — next slide.",
         size=14, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 13b: UC4/UC5 — the real problem and the fix
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "UC4 and UC5: the problem isn't speed, it's waste",
         size=26, bold=True)

# UC4
add_pill(slide, 0.8, 1.4, 3.5, "UC4: Artifact metadata", color=PURPLE)
add_text(slide, 0.8, 2.0, 5.5, 2.5,
         "Today: Client API is already fast (5 calls, ~50ms). "
         "The metadata service returns artifact name, sha, type, "
         "location without loading data.\n\n"
         "Problem: no way to list artifacts across tasks or runs. "
         "\"Find the task where accuracy was highest\" requires iterating "
         "every task and loading the artifact.\n\n"
         "Fix: Add a cross-task artifact listing endpoint to the metadata "
         "service. One call returns artifact metadata for all tasks in a step, "
         "filterable by name.",
         size=14)

# UC5
add_pill(slide, 6.8, 1.4, 3.5, "UC5: Bounded logs", color=PURPLE)
add_text(slide, 6.8, 2.0, 5.5, 2.5,
         "Today: task.stdout loads the ENTIRE log as a string. "
         "A training step with 50MB of output? Agent gets 50MB "
         "when it only needs the last 10 lines.\n\n"
         "Problem: there's no \"tail\" or \"grep\" for logs. "
         "The metadata service has no log endpoint at all. "
         "The UI Backend has one but with cache overhead.\n\n"
         "Fix: Two options:\n"
         "1. Client-side: read log file, seek to end, return last N lines\n"
         "2. Service-side: add /logs endpoint with ?_limit=N to metadata service\n"
         "Either way, agent gets bounded output.",
         size=14)

# Bottom summary
add_text(slide, 0.8, 5.5, 11.5, 1.5,
         "The insight: UC4/UC5 aren't about fewer calls. They're about bounded data.\n"
         "An agent shouldn't need to download 50MB of logs to read 10 lines,\n"
         "or iterate 1000 tasks to find one artifact name.",
         size=16, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 14: Section divider - Solution
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1.5, 2.8, 10, 2,
         "What we're building",
         size=40, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 15: Three layers
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "Three layers, each independently useful",
         size=28, bold=True)

layers = [
    ("Layer 1", "Extension Package", "metaflow-agent-client",
     "Smart functions using filter_tasks_by_metadata + existing API",
     "No core changes. Works today on service >= 2.5.0", GREEN),
    ("Layer 2", "Metadata Service", "~150 lines",
     "Add _limit, _order to endpoints (DB already supports it)",
     "Benefits everyone. Small, targeted changes.", PURPLE),
    ("Layer 3", "Client API", "ServiceMetadataProvider",
     "Automatically use new query params (version-gated)",
     "Transparent speedup. Existing code gets faster.", BLUE),
]

for i, (label, title, subtitle, desc, note, color) in enumerate(layers):
    y = 1.5 + i * 1.8
    add_pill(slide, 0.8, y, 1.8, label, color=color)
    add_text(slide, 3.0, y - 0.05, 4, 0.5, title, size=22, bold=True)
    add_text(slide, 3.0, y + 0.4, 4, 0.3, subtitle, size=14, color=TEXT_LIGHT)
    add_text(slide, 7.5, y, 5, 0.5, desc, size=15)
    add_text(slide, 7.5, y + 0.5, 5, 0.4, note, size=13, color=color, bold=True)

# Key point
add_text(slide, 0.8, 6.6, 11.5, 0.5,
         "No dependency on UI Backend. Every deployment gets the improvement.",
         size=18, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 16: Before/After code
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "Before and After",
         size=28, bold=True)

add_text(slide, 0.8, 1.3, 5.5, 0.4,
         "Today: 56 HTTP calls", size=18, bold=True, color=RED)
add_code_block(slide, 0.8, 1.8, 5.5, 3.0, """# Agent iterates every task
failures = []
for step in run:
    for task in step:
        if not task.successful:
            failures.append(task.pathspec)

# 56 HTTP calls, ~3800ms
# 50 artifact fetches
# 50 datastore reads
# 50 unpickle operations""", size=13)

add_text(slide, 7.0, 1.3, 5.5, 0.4,
         "After: 4 HTTP calls", size=18, bold=True, color=GREEN)
add_code_block(slide, 7.0, 1.8, 5.5, 3.0, """# Agent calls one function
from metaflow_extensions.agent_client import (
    find_failures
)

result = find_failures("ForeachFlow/8")
# result["failures"]
#   -> ["ForeachFlow/8/process/126"]

# 4 HTTP calls, ~350ms
# 0 artifact fetches
# 0 datastore reads""", size=13)

add_text(slide, 0.8, 5.3, 11.5, 1.0,
         "Same answer. 93% fewer calls. 8.6x faster.\n"
         "Works on every deployment with metadata service >= 2.5.0.",
         size=20, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 17: Timeline
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "GSoC Timeline: 350 hours, 16 weeks",
         size=28, bold=True)

phases = [
    ("Phase 1", "Weeks 1-4", "90h", "Foundation",
     "Extension package + 6 utility functions + tests",
     GREEN),
    ("Phase 2", "Weeks 5-8", "90h", "Core Implementation",
     "Metadata service _limit/_order + provider enhancement",
     PURPLE),
    ("Phase 3", "Weeks 9-12", "90h", "Agent Integration",
     "Agent simulation, tool schemas, robustness testing",
     BLUE),
    ("Phase 4", "Weeks 13-16", "80h", "Documentation & Polish",
     "User/agent docs, PR preparation, final report",
     TEXT_LIGHT),
]

for i, (phase, weeks, hours, title, desc, color) in enumerate(phases):
    y = 1.5 + i * 1.4
    # Phase pill
    add_pill(slide, 0.8, y, 1.8, phase, color=color)
    # Timeline bar
    bar_width = 8.5
    add_rect(slide, 3.0, y + 0.05, bar_width, 0.35, color)
    add_text(slide, 3.1, y + 0.02, 2, 0.35, f"{weeks} | {hours}",
             size=12, bold=True, color=WHITE)
    add_text(slide, 5.5, y + 0.02, 6, 0.35, title,
             size=14, bold=True, color=WHITE)
    # Description below
    add_text(slide, 3.0, y + 0.5, 9, 0.4, desc, size=14, color=TEXT_LIGHT)

add_text(slide, 0.8, 6.8, 11.5, 0.4,
         "Each phase delivers standalone, testable, demonstrable value.",
         size=16, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 18: Why this matters
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.6,
         "Why this matters",
         size=28, bold=True)

categories = [
    ("For agents", [
        "Debug flows without hammering the backend",
        "Structured responses, bounded queries, predictable cost",
        "Tool schemas for MCP/function-calling integration",
    ], PURPLE),
    ("For all users", [
        "Pagination and ordering on metadata service (benefits everyone)",
        "Faster Client API iteration (transparent, version-gated)",
        "Audit reveals where to invest next in the API",
    ], GREEN),
    ("For Metaflow", [
        "First systematic mapping of Client API to HTTP calls",
        "Extension package pattern for community contributions",
        "Foundation for richer programmatic access",
    ], BLUE),
]

for i, (title, points, color) in enumerate(categories):
    x = 0.8 + i * 4.0
    add_pill(slide, x, 1.5, 3.2, title, color=color)
    for j, point in enumerate(points):
        add_text(slide, x, 2.2 + j * 0.8, 3.5, 0.7,
                 point, size=14, color=TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 19: Live demo placeholder
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1.5, 2.5, 10, 2,
         "Live Demo",
         size=48, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.0, 10, 1.5,
         "Claude debugging a Metaflow flow\n"
         "Side-by-side: naive path vs smart path\n"
         "Real HTTP calls, real Metaflow UI, real numbers",
         size=20, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 6.0, 10, 0.5,
         "[YouTube video link / screen recording]",
         size=16, color=PURPLE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 20: Thank you
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1.5, 2.0, 10, 1.5,
         "Thank you",
         size=48, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.5, 10, 0.5,
         "landigf | MSc CS, ETH Zurich | EASL Lab",
         size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.2, 10, 0.5,
         "GSoC 2026 | Mentor: Valay Dave | Outerbounds",
         size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_multiline(slide, 3.5, 5.2, 6, 1.5, [
    "Code: github.com/landigf/metaflow-agent-client",
    "Benchmarks: demo/benchmark_three_paths.py",
    "RFC: docs/RFC_agent_friendly_client.md",
    "Full audit: demo/full_audit.py (all 6 use cases)",
], size=14, color=TEXT_LIGHT)

# ── Save ──
output_path = "/Users/landigf/Desktop/Code/GSoC/docs/GSoC_2026_Agent_Friendly_Metaflow.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
